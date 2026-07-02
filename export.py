"""
Exportación a Excel y PowerPoint con formato ejecutivo.
Excel: libro con 5 hojas (Resumen, Mensual, Trimestral, Semestral, Anual).
PPTX:  presentación ejecutiva con KPIs, tabla anual y gráficos.
"""
import io
from datetime import datetime

from openpyxl import Workbook
from openpyxl.styles import (
    PatternFill, Font, Alignment, Border, Side, numbers as xl_numbers
)
from openpyxl.utils import get_column_letter
import pandas as pd

from model import ProjectionParams, aggregate_by_period

# ── Estilos ───────────────────────────────────────────────────────────────────

F_DARK       = PatternFill("solid", fgColor="0F2037")
F_NAVY       = PatternFill("solid", fgColor="1E3A5F")
F_ALT        = PatternFill("solid", fgColor="0A1628")
F_POS        = PatternFill("solid", fgColor="052E16")
F_NEG        = PatternFill("solid", fgColor="450A0A")

FONT_TITLE   = Font(name="Calibri", color="60A5FA", size=16, bold=True)
FONT_HDR     = Font(name="Calibri", color="FFFFFF", size=10, bold=True)
FONT_SUB     = Font(name="Calibri", color="60A5FA", size=10, bold=True)
FONT_NORMAL  = Font(name="Calibri", color="E2E8F0", size=10)
FONT_POS     = Font(name="Calibri", color="4ADE80", size=10, bold=True)
FONT_NEG     = Font(name="Calibri", color="F87171", size=10, bold=True)
FONT_LABEL   = Font(name="Calibri", color="94A3B8", size=10)

ALIGN_CENTER = Alignment(horizontal="center", vertical="center", wrap_text=True)
ALIGN_RIGHT  = Alignment(horizontal="right",  vertical="center")
ALIGN_LEFT   = Alignment(horizontal="left",   vertical="center")

BORDER_THIN  = Border(
    bottom=Side(border_style="thin", color="1E3A5F"),
    top=Side(border_style="thin", color="1E3A5F"),
)

FMT_CURR = '"$"#,##0.0'
FMT_PCT  = '#,##0.00"%"'
FMT_INT  = '0'
FMT_NONE = '#,##0.0'

COL_FORMATS = {
    "Mes": FMT_INT, "Año": FMT_INT, "Mes_en_año": FMT_INT,
    "Tasa_colocacion_pct": FMT_PCT, "Costo_fondo_pct": FMT_PCT,
    "Margen_neto_pct": FMT_PCT, "Eficiencia_pct": FMT_PCT, "ROA_pct": FMT_PCT,
}

COL_NAMES_ES = {
    "Mes": "Mes", "Año": "Año", "Periodo": "Período",
    "Cartera": "Cartera (M$)", "Crecimiento_mensual": "Crecimiento (M$)",
    "Tasa_colocacion_pct": "Tasa Coloc. (%)", "Costo_fondo_pct": "Costo Fdo. (%)",
    "Ingresos_factoring": "Ingresos (M$)", "Costo_fondo": "Costo Fondo (M$)",
    "Margen_financiero": "Margen Fin. (M$)", "Costos_operacionales": "Op. (M$)",
    "Remuneraciones": "Rem. (M$)", "Otros_gastos": "Otros (M$)",
    "Total_costos": "Total Costos (M$)", "Resultado_neto": "Resultado Neto (M$)",
    "Margen_neto_pct": "Margen Neto (%)", "Eficiencia_pct": "Eficiencia (%)",
    "ROA_pct": "ROA (%)",
}


def _write_df(ws, df: pd.DataFrame, start_row: int = 1, freeze: bool = True):
    """Escribe DataFrame con cabecera estilizada."""
    ws.sheet_view.showGridLines = False

    # Cabecera
    for ci, col in enumerate(df.columns, start=1):
        cell = ws.cell(row=start_row, column=ci, value=COL_NAMES_ES.get(col, col))
        cell.fill      = F_DARK
        cell.font      = FONT_HDR
        cell.alignment = ALIGN_CENTER
        cell.border    = BORDER_THIN

    # Datos
    for ri, row_data in enumerate(df.itertuples(index=False), start=start_row + 1):
        alt = (ri % 2 == 0)
        for ci, value in enumerate(row_data, start=1):
            cell = ws.cell(row=ri, column=ci, value=value)
            cell.font      = FONT_NORMAL
            cell.alignment = ALIGN_RIGHT

            col_name = df.columns[ci - 1]

            if alt:
                cell.fill = F_ALT

            # Formato numérico
            if col_name in COL_FORMATS:
                cell.number_format = COL_FORMATS[col_name]
            elif col_name in ("Periodo",):
                cell.alignment = ALIGN_LEFT
                cell.font = FONT_NORMAL
            elif isinstance(value, float):
                cell.number_format = FMT_CURR

            # Colorear resultado neto
            if col_name == "Resultado_neto" and isinstance(value, (int, float)):
                cell.fill = F_POS if value >= 0 else F_NEG
                cell.font = FONT_POS if value >= 0 else FONT_NEG

    # Congelar fila de cabecera
    if freeze:
        ws.freeze_panes = ws.cell(row=start_row + 1, column=1)

    # Ajustar anchos
    for ci, col in enumerate(df.columns, start=1):
        letter = get_column_letter(ci)
        header_len = len(COL_NAMES_ES.get(col, col))
        max_data   = max(
            (len(str(v)) for v in df[col] if v is not None),
            default=0,
        )
        ws.column_dimensions[letter].width = min(max(header_len, max_data) + 3, 22)

    # Altura de cabecera
    ws.row_dimensions[start_row].height = 32


def export_to_excel(df_monthly: pd.DataFrame, params: ProjectionParams) -> bytes:
    """Genera workbook Excel y retorna bytes para descarga."""
    wb = Workbook()

    # ── Hoja 1: Resumen Ejecutivo ─────────────────────────────────────────────
    ws_sum = wb.active
    ws_sum.title = "Resumen Ejecutivo"
    ws_sum.sheet_view.showGridLines = False
    ws_sum.column_dimensions["A"].width = 3
    ws_sum.column_dimensions["B"].width = 32
    ws_sum.column_dimensions["C"].width = 22

    ws_sum.merge_cells("B2:F2")
    ws_sum["B2"] = "PROYECCIÓN FINANCIERA — FACTORING"
    ws_sum["B2"].font      = FONT_TITLE
    ws_sum["B2"].fill      = F_DARK
    ws_sum["B2"].alignment = ALIGN_CENTER

    ws_sum["B3"] = f"Generado: {datetime.now().strftime('%d/%m/%Y %H:%M')}  |  Horizonte: 6 años (72 meses)"
    ws_sum["B3"].font = FONT_LABEL

    # Parámetros base
    ws_sum["B5"] = "PARÁMETROS BASE"
    ws_sum["B5"].font = FONT_SUB
    ws_sum["B5"].fill = F_NAVY

    param_rows = [
        ("Cartera Inicial",              f"${params.initial_portfolio:,.0f}M"),
        ("Crecimiento Anual de Cartera", f"${params.annual_portfolio_growth:,.0f}M"),
        ("Tasa de Colocación (mensual)", f"{params.placement_rate:.2f}%"),
        ("Costo de Fondo (mensual)",     f"{params.funding_cost_rate:.2f}%"),
        ("Costos Operacionales Base",    f"${params.current_op_costs:,.1f}M/mes"),
        ("Remuneraciones Base",          f"${params.current_remuneration:,.1f}M/mes"),
        ("Incremento Anual Op.",         f"${params.annual_op_increment:,.1f}M/año"),
        ("Otros Gastos",                 f"${params.other_expenses:,.1f}M/mes"),
    ]
    if params.overrides:
        param_rows.append(("Modificaciones Programadas", f"{len(params.overrides)} override(s)"))

    for i, (label, value) in enumerate(param_rows, start=6):
        ws_sum[f"B{i}"] = label
        ws_sum[f"C{i}"] = value
        ws_sum[f"B{i}"].font = FONT_LABEL
        ws_sum[f"C{i}"].font = Font(name="Calibri", color="FFFFFF", size=10, bold=True)
        ws_sum[f"B{i}"].alignment = ALIGN_LEFT
        ws_sum[f"C{i}"].alignment = ALIGN_LEFT
        if i % 2 == 0:
            ws_sum[f"B{i}"].fill = F_ALT
            ws_sum[f"C{i}"].fill = F_ALT

    # KPIs finales
    last = df_monthly.iloc[-1]
    total_net  = df_monthly["Resultado_neto"].sum()
    avg_margin = df_monthly["Margen_neto_pct"].mean()
    final_portfolio = last["Cartera"]

    kpi_start = 18
    ws_sum[f"B{kpi_start}"] = "KPIs PROYECTADOS (6 AÑOS)"
    ws_sum[f"B{kpi_start}"].font = FONT_SUB
    ws_sum[f"B{kpi_start}"].fill = F_NAVY

    kpi_rows = [
        ("Cartera Final (Mes 72)",        f"${final_portfolio:,.0f}M"),
        ("Resultado Neto Total",           f"${total_net:,.0f}M"),
        ("Margen Neto Promedio",           f"{avg_margin:.1f}%"),
        ("ROA Final (Mes 72)",             f"{last['ROA_pct']:.2f}%"),
        ("Eficiencia Final (Mes 72)",      f"{last['Eficiencia_pct']:.1f}%"),
    ]
    for i, (label, value) in enumerate(kpi_rows, start=kpi_start + 1):
        ws_sum[f"B{i}"] = label
        ws_sum[f"C{i}"] = value
        ws_sum[f"B{i}"].font = FONT_LABEL
        ws_sum[f"C{i}"].font = FONT_POS if "M" in str(value) or "%" in str(value) else FONT_NORMAL
        if i % 2 == 0:
            ws_sum[f"B{i}"].fill = F_ALT
            ws_sum[f"C{i}"].fill = F_ALT

    # Tabla anual en hoja resumen
    df_annual = aggregate_by_period(df_monthly, "Anual")
    sum_cols  = ["Periodo", "Cartera", "Ingresos_factoring", "Margen_financiero",
                 "Total_costos", "Resultado_neto", "Margen_neto_pct", "ROA_pct"]
    df_a_disp = df_annual[[c for c in sum_cols if c in df_annual.columns]]

    ws_sum[f"B{kpi_start + 9}"] = "RESUMEN ANUAL"
    ws_sum[f"B{kpi_start + 9}"].font = FONT_SUB
    ws_sum[f"B{kpi_start + 9}"].fill = F_NAVY

    _write_df(ws_sum, df_a_disp, start_row=kpi_start + 10, freeze=False)

    # ── Hojas de proyección ───────────────────────────────────────────────────
    periods = [
        ("Proyección Mensual",    "Mensual"),
        ("Trimestral",            "Trimestral"),
        ("Semestral",             "Semestral"),
        ("Anual",                 "Anual"),
    ]

    monthly_cols = ["Mes", "Año", "Cartera", "Ingresos_factoring", "Costo_fondo",
                    "Margen_financiero", "Costos_operacionales", "Remuneraciones",
                    "Otros_gastos", "Total_costos", "Resultado_neto",
                    "Margen_neto_pct", "Eficiencia_pct", "ROA_pct"]

    period_cols = ["Periodo", "Cartera", "Ingresos_factoring", "Costo_fondo",
                   "Margen_financiero", "Costos_operacionales", "Remuneraciones",
                   "Total_costos", "Resultado_neto", "Margen_neto_pct", "Eficiencia_pct", "ROA_pct"]

    for sheet_name, period in periods:
        ws = wb.create_sheet(sheet_name)
        df_p = aggregate_by_period(df_monthly, period)

        if period == "Mensual":
            cols = [c for c in monthly_cols if c in df_p.columns]
        else:
            cols = [c for c in period_cols if c in df_p.columns]

        _write_df(ws, df_p[cols])

    # ── Serializar a bytes ────────────────────────────────────────────────────
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.read()


# ── PowerPoint export ─────────────────────────────────────────────────────────

def export_to_pptx(df_monthly: pd.DataFrame, params: ProjectionParams,
                   scenario_name: str = "") -> bytes:
    """Genera presentación ejecutiva PowerPoint y retorna bytes para descarga."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.oxml.ns import qn
    from lxml import etree

    # Paleta SMB
    C_DARK   = RGBColor(0x0F, 0x20, 0x37)
    C_NAVY   = RGBColor(0x1E, 0x3A, 0x5F)
    C_ORANGE = RGBColor(0xF4, 0x79, 0x20)
    C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    C_GRAY   = RGBColor(0xA8, 0xA8, 0xA8)
    C_GREEN  = RGBColor(0x4A, 0xDE, 0x80)
    C_RED    = RGBColor(0xF8, 0x71, 0x71)

    W = Inches(13.33)
    H = Inches(7.5)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    df_annual = aggregate_by_period(df_monthly, "Anual")
    last      = df_monthly.iloc[-1]
    m1        = df_monthly.iloc[0]
    total_net = df_monthly["Resultado_neto"].sum()

    # ── Helpers ───────────────────────────────────────────────────────────────

    def fill_slide(slide, color: RGBColor):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_rect(slide, left, top, width, height, color: RGBColor, alpha=None):
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_text(slide, text, left, top, width, height,
                 size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
                 italic=False, wrap=True):
        txb = slide.shapes.add_textbox(left, top, width, height)
        tf  = txb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color or C_WHITE
        return txb

    def kpi_card(slide, left, top, w, h, label, value, sublabel="", pos=True):
        add_rect(slide, left, top, w, h, C_NAVY)
        vclr = C_GREEN if pos else C_RED
        add_text(slide, label,    left+Inches(0.1), top+Inches(0.05),  w-Inches(0.2), Inches(0.35),
                 size=9, color=C_GRAY, align=PP_ALIGN.CENTER)
        add_text(slide, value,    left+Inches(0.1), top+Inches(0.35),  w-Inches(0.2), Inches(0.6),
                 size=22, bold=True, color=vclr, align=PP_ALIGN.CENTER)
        if sublabel:
            add_text(slide, sublabel, left+Inches(0.1), top+Inches(0.9),  w-Inches(0.2), Inches(0.25),
                     size=8, color=C_GRAY, align=PP_ALIGN.CENTER)

    # ── Slide 1: Portada ──────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    fill_slide(s1, C_DARK)
    add_rect(s1, 0, 0, W, Inches(0.08), C_ORANGE)
    add_rect(s1, 0, H - Inches(0.08), W, Inches(0.08), C_ORANGE)
    add_rect(s1, Inches(0.5), Inches(1.8), Inches(0.06), Inches(4), C_ORANGE)

    title = scenario_name or "Proyección Financiera"
    add_text(s1, "SMB SERVICIOS FINANCIEROS S.A.",
             Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.8),
             size=16, bold=True, color=C_ORANGE)
    add_text(s1, title,
             Inches(0.8), Inches(2.7), Inches(11.5), Inches(1.2),
             size=38, bold=True, color=C_WHITE)
    add_text(s1, "Horizonte 6 años (72 meses)  ·  Valores en millones de pesos (M$)",
             Inches(0.8), Inches(3.9), Inches(11.5), Inches(0.5),
             size=14, color=C_GRAY)
    add_text(s1, f"Generado: {datetime.now().strftime('%d/%m/%Y')}",
             Inches(0.8), Inches(4.6), Inches(8), Inches(0.4),
             size=12, color=C_GRAY)

    # ── Slide 2: Parámetros base ──────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    fill_slide(s2, C_DARK)
    add_rect(s2, 0, 0, W, Inches(1.1), C_NAVY)
    add_text(s2, "PARÁMETROS DEL MODELO", Inches(0.4), Inches(0.25), Inches(10), Inches(0.6),
             size=22, bold=True, color=C_ORANGE)
    add_text(s2, "Calibrado con EERR Mayo 2026", Inches(0.4), Inches(0.75), Inches(8), Inches(0.4),
             size=11, color=C_GRAY)

    param_data = [
        ("Cartera Factoring Inicial",     f"M${params.initial_portfolio:,.0f}"),
        ("Crecimiento Anual Factoring",   f"M${params.annual_portfolio_growth:,.0f}/año"),
        ("Tasa de Colocación",            f"{params.placement_rate:.2f}%/mes"),
        ("Costo de Fondo (variable)",     f"{params.funding_cost_rate:.2f}%/mes"),
        ("Costo Fijo de Fondo",          f"M${params.base_funding_cost_mm:.1f}/mes"),
        ("Cartera Leasing Inicial (UF)", f"{params.leasing_portfolio_uf:,.0f} UF"),
        ("Tasa Anual Leasing",           f"{params.leasing_annual_rate:.1f}%"),
        ("Crecimiento Leasing",          f"M${params.leasing_annual_growth_mm:,.0f}/año"),
        ("Remuneraciones Base",          f"M${params.current_remuneration:,.1f}/mes"),
        ("Costos Operacionales Base",    f"M${params.current_op_costs:,.1f}/mes"),
        ("Otros Gastos",                 f"M${params.other_expenses:,.1f}/mes"),
        ("Tasa de Provisión",            f"{params.provision_rate:.1f}% anual"),
        ("Modificaciones programadas",   f"{len(params.overrides)} override(s)"),
    ]

    cols2 = 2
    per_col = len(param_data) // cols2 + 1
    col_w   = Inches(6.2)
    row_h   = Inches(0.4)
    for i, (lbl, val) in enumerate(param_data):
        col_i = i // per_col
        row_i = i % per_col
        lx = Inches(0.4) + col_i * col_w
        ty = Inches(1.25) + row_i * row_h
        bg = C_NAVY if i % 2 == 0 else C_DARK
        add_rect(s2, lx, ty, col_w - Inches(0.1), row_h - Inches(0.04), bg)
        add_text(s2, lbl, lx + Inches(0.1), ty + Inches(0.06),
                 Inches(3.5), row_h, size=10, color=C_GRAY)
        add_text(s2, val, lx + Inches(3.6), ty + Inches(0.06),
                 Inches(2.5), row_h, size=10, bold=True, color=C_WHITE, align=PP_ALIGN.RIGHT)

    # ── Slide 3: KPIs clave ───────────────────────────────────────────────────
    s3 = prs.slides.add_slide(blank)
    fill_slide(s3, C_DARK)
    add_rect(s3, 0, 0, W, Inches(1.1), C_NAVY)
    add_text(s3, "KPIs PROYECTADOS", Inches(0.4), Inches(0.25), Inches(10), Inches(0.6),
             size=22, bold=True, color=C_ORANGE)
    add_text(s3, "Resultado a 6 años · Mes 1 vs Mes 72", Inches(0.4), Inches(0.75), Inches(8), Inches(0.4),
             size=11, color=C_GRAY)

    kpi_w = Inches(2.45)
    kpi_h = Inches(1.35)
    kpi_gap = Inches(0.18)
    kpi_top1 = Inches(1.3)
    kpi_top2 = kpi_top1 + kpi_h + kpi_gap

    kpis_row1 = [
        ("Cartera Final (Mes 72)",  f"M${last['Cartera']:,.0f}",       "", True),
        ("Resultado Neto Total",    f"M${total_net:,.0f}",              "6 años acumulado", total_net >= 0),
        ("Margen Neto Prom.",       f"{df_monthly['Margen_neto_pct'].mean():.1f}%", "promedio 72 meses", True),
        ("Cartera Leasing Final",   f"M${last['Cartera_leasing']:,.0f}", "Mes 72", True),
        ("Imp. Renta Total",        f"M${df_monthly['Impuesto'].sum():,.0f}", "acumulado 6 años", True),
    ]
    kpis_row2 = [
        ("Resultado Mes 1",         f"M${m1['Resultado_neto']:,.1f}",   "base EERR May-26", m1["Resultado_neto"] >= 0),
        ("Resultado Mes 72",        f"M${last['Resultado_neto']:,.1f}", "estado estable",   last["Resultado_neto"] >= 0),
        ("ROA Final",               f"{last['ROA_pct']:.2f}%",          "anualizado Mes 72", last['ROA_pct'] >= 0),
        ("Eficiencia Final",        f"{last['Eficiencia_pct']:.1f}%",   "costos / margen Mes 72", True),
        ("Patrimonio Final",        f"M${last['Patrimonio']:,.0f}",      "Mes 72", True),
    ]

    for row_i, (kpis, top) in enumerate([(kpis_row1, kpi_top1), (kpis_row2, kpi_top2)]):
        for j, (lbl, val, sub, pos) in enumerate(kpis):
            lx = Inches(0.4) + j * (kpi_w + kpi_gap)
            kpi_card(s3, lx, top, kpi_w, kpi_h, lbl, val, sub, pos)

    # ── Slide 4: Resultados anuales ───────────────────────────────────────────
    s4 = prs.slides.add_slide(blank)
    fill_slide(s4, C_DARK)
    add_rect(s4, 0, 0, W, Inches(1.0), C_NAVY)
    add_text(s4, "RESUMEN ANUAL P&L", Inches(0.4), Inches(0.2), Inches(10), Inches(0.6),
             size=22, bold=True, color=C_ORANGE)
    add_text(s4, "Valores en M$ (millones de pesos)", Inches(0.4), Inches(0.72), Inches(8), Inches(0.35),
             size=10, color=C_GRAY)

    tbl_cols = ["Periodo", "Ingresos_factoring", "Ingresos_leasing",
                "Costo_fondo", "Margen_financiero",
                "Total_costos", "Resultado_neto", "Margen_neto_pct"]
    tbl_hdrs = ["Período", "Ing. Fact.", "Ing. Leas.", "Costo Fondo",
                "Margen", "Tot. Costos", "Res. Neto", "Margen %"]

    rows_tbl = len(df_annual) + 1  # +1 cabecera
    cols_tbl = len(tbl_cols)
    tbl_left = Inches(0.3)
    tbl_top  = Inches(1.15)
    tbl_w    = W - Inches(0.6)
    tbl_h    = H - Inches(1.35)

    table = s4.shapes.add_table(rows_tbl, cols_tbl, tbl_left, tbl_top, tbl_w, tbl_h).table

    col_widths = [Inches(1.3)] + [Inches(1.43)] * (cols_tbl - 1)
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw

    for ci, hdr in enumerate(tbl_hdrs):
        cell = table.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_DARK
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.bold  = True
        run.font.color.rgb = C_ORANGE
        run.font.size  = Pt(9)

    for ri, (_, row) in enumerate(df_annual.iterrows(), start=1):
        bg = C_NAVY if ri % 2 == 0 else C_DARK
        for ci, col in enumerate(tbl_cols):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            val = row[col]
            if col == "Periodo":
                txt = str(val)
                clr = C_WHITE
                al  = PP_ALIGN.LEFT
            elif col == "Margen_neto_pct":
                txt = f"{val:.1f}%"
                clr = C_WHITE
                al  = PP_ALIGN.CENTER
            else:
                txt = f"${val:,.0f}"
                clr = (C_GREEN if val >= 0 else C_RED) if col == "Resultado_neto" else C_WHITE
                al  = PP_ALIGN.RIGHT
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = al
            p.text = txt
            run = p.runs[0]
            run.font.size  = Pt(9)
            run.font.color.rgb = clr
            run.font.bold  = (col == "Resultado_neto")

    # ── Slide 5: Gráfico Resultado Neto Anual ────────────────────────────────
    s5 = prs.slides.add_slide(blank)
    fill_slide(s5, C_DARK)
    add_rect(s5, 0, 0, W, Inches(1.0), C_NAVY)
    add_text(s5, "RESULTADO NETO POR AÑO", Inches(0.4), Inches(0.2), Inches(10), Inches(0.6),
             size=22, bold=True, color=C_ORANGE)
    add_text(s5, "Millones de pesos (M$)", Inches(0.4), Inches(0.72), Inches(8), Inches(0.35),
             size=10, color=C_GRAY)

    cd5 = ChartData()
    cd5.categories = [str(r["Periodo"]) for _, r in df_annual.iterrows()]
    cd5.add_series("Resultado Neto (M$)", [round(r["Resultado_neto"], 1) for _, r in df_annual.iterrows()])

    chart5 = s5.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.0), cd5
    ).chart

    chart5.has_legend = False
    chart5.has_title  = False
    chart5.plots[0].series[0].format.fill.solid()
    chart5.plots[0].series[0].format.fill.fore_color.rgb = C_ORANGE

    sv5 = chart5.value_axis
    sv5.has_major_gridlines = True
    sv5.major_gridlines.format.line.color.rgb = RGBColor(0x2D, 0x2D, 0x2D)

    # ── Slide 6: Gráfico Cartera ──────────────────────────────────────────────
    s6 = prs.slides.add_slide(blank)
    fill_slide(s6, C_DARK)
    add_rect(s6, 0, 0, W, Inches(1.0), C_NAVY)
    add_text(s6, "EVOLUCIÓN DE CARTERA", Inches(0.4), Inches(0.2), Inches(10), Inches(0.6),
             size=22, bold=True, color=C_ORANGE)
    add_text(s6, "Factoring + Leasing en M$", Inches(0.4), Inches(0.72), Inches(8), Inches(0.35),
             size=10, color=C_GRAY)

    cd6 = ChartData()
    cd6.categories = [str(r["Periodo"]) for _, r in df_annual.iterrows()]
    cd6.add_series("Cartera Factoring",  [round(r["Cartera"], 0) for _, r in df_annual.iterrows()])
    cd6.add_series("Cartera Leasing",    [round(r["Cartera_leasing"], 0) for _, r in df_annual.iterrows()])

    chart6 = s6.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.0), cd6
    ).chart

    chart6.has_legend = True
    chart6.has_title  = False
    chart6.plots[0].series[0].format.fill.solid()
    chart6.plots[0].series[0].format.fill.fore_color.rgb = C_ORANGE
    chart6.plots[0].series[1].format.fill.solid()
    chart6.plots[0].series[1].format.fill.fore_color.rgb = RGBColor(0xA7, 0x8B, 0xFA)

    # ── Slide 7: Margen e Ingresos ────────────────────────────────────────────
    s7 = prs.slides.add_slide(blank)
    fill_slide(s7, C_DARK)
    add_rect(s7, 0, 0, W, Inches(1.0), C_NAVY)
    add_text(s7, "INGRESOS, COSTOS Y MARGEN", Inches(0.4), Inches(0.2), Inches(10), Inches(0.6),
             size=22, bold=True, color=C_ORANGE)
    add_text(s7, "Millones de pesos anuales (M$)", Inches(0.4), Inches(0.72), Inches(8), Inches(0.35),
             size=10, color=C_GRAY)

    cd7 = ChartData()
    cd7.categories = [str(r["Periodo"]) for _, r in df_annual.iterrows()]
    cd7.add_series("Ingresos Totales",  [round(r["Ingresos_total"], 0) for _, r in df_annual.iterrows()])
    cd7.add_series("Total Costos",      [round(abs(r["Total_costos"]), 0) for _, r in df_annual.iterrows()])
    cd7.add_series("Margen Financiero", [round(r["Margen_financiero"], 0) for _, r in df_annual.iterrows()])

    chart7 = s7.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.0), cd7
    ).chart

    chart7.has_legend = True
    chart7.has_title  = False
    colors7 = [C_ORANGE, C_RED, RGBColor(0xC8, 0xC8, 0xC8)]
    for si, clr in enumerate(colors7):
        chart7.plots[0].series[si].format.line.color.rgb = clr
        chart7.plots[0].series[si].format.line.width = Pt(2.5)

    # ── Serializar ────────────────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
