"""
Exportación a Excel y PowerPoint con formato ejecutivo.
Excel: libro con 5 hojas (Resumen, Mensual, Trimestral, Semestral, Anual).
PPTX:  presentación ejecutiva con KPIs, tabla anual y gráficos.
"""
import io
from datetime import datetime

from openpyxl import Workbook
from openpyxl.styles import (
    PatternFill, Font, Alignment, Border, Side, numbers as xl_numbers
)
from openpyxl.utils import get_column_letter
import pandas as pd

from model import ProjectionParams, aggregate_by_period

# ── Estilos ───────────────────────────────────────────────────────────────────

F_DARK       = PatternFill("solid", fgColor="0F2037")
F_NAVY       = PatternFill("solid", fgColor="1E3A5F")
F_ALT        = PatternFill("solid", fgColor="0A1628")
F_POS        = PatternFill("solid", fgColor="052E16")
F_NEG        = PatternFill("solid", fgColor="450A0A")

FONT_TITLE   = Font(name="Calibri", color="60A5FA", size=16, bold=True)
FONT_HDR     = Font(name="Calibri", color="FFFFFF", size=10, bold=True)
FONT_SUB     = Font(name="Calibri", color="60A5FA", size=10, bold=True)
FONT_NORMAL  = Font(name="Calibri", color="E2E8F0", size=10)
FONT_POS     = Font(name="Calibri", color="4ADE80", size=10, bold=True)
FONT_NEG     = Font(name="Calibri", color="F87171", size=10, bold=True)
FONT_LABEL   = Font(name="Calibri", color="94A3B8", size=10)

ALIGN_CENTER = Alignment(horizontal="center", vertical="center", wrap_text=True)
ALIGN_RIGHT  = Alignment(horizontal="right",  vertical="center")
ALIGN_LEFT   = Alignment(horizontal="left",   vertical="center")

BORDER_THIN  = Border(
    bottom=Side(border_style="thin", color="1E3A5F"),
    top=Side(border_style="thin", color="1E3A5F"),
)

FMT_CURR = '"$"#,##0.0'
FMT_PCT  = '#,##0.00"%"'
FMT_INT  = '0'
FMT_NONE = '#,##0.0'

COL_FORMATS = {
    "Mes": FMT_INT, "Año": FMT_INT, "Mes_en_año": FMT_INT,
    "Tasa_colocacion_pct": FMT_PCT, "Costo_fondo_pct": FMT_PCT,
    "Margen_neto_pct": FMT_PCT, "Eficiencia_pct": FMT_PCT, "ROA_pct": FMT_PCT,
}

COL_NAMES_ES = {
    "Mes": "Mes", "Año": "Año", "Periodo": "Período",
    "Cartera": "Cartera (M$)", "Crecimiento_mensual": "Crecimiento (M$)",
    "Tasa_colocacion_pct": "Tasa Coloc. (%)", "Costo_fondo_pct": "Costo Fdo. (%)",
    "Ingresos_factoring": "Ingresos (M$)", "Costo_fondo": "Costo Fondo (M$)",
    "Margen_financiero": "Margen Fin. (M$)", "Costos_operacionales": "Op. (M$)",
    "Remuneraciones": "Rem. (M$)", "Otros_gastos": "Otros (M$)",
    "Total_costos": "Total Costos (M$)", "Resultado_neto": "Resultado Neto (M$)",
    "Margen_neto_pct": "Margen Neto (%)", "Eficiencia_pct": "Eficiencia (%)",
    "ROA_pct": "ROA (%)",
}


def _write_df(ws, df: pd.DataFrame, start_row: int = 1, freeze: bool = True):
    """Escribe DataFrame con cabecera estilizada."""
    ws.sheet_view.showGridLines = False

    # Cabecera
    for ci, col in enumerate(df.columns, start=1):
        cell = ws.cell(row=start_row, column=ci, value=COL_NAMES_ES.get(col, col))
        cell.fill      = F_DARK
        cell.font      = FONT_HDR
        cell.alignment = ALIGN_CENTER
        cell.border    = BORDER_THIN

    # Datos
    for ri, row_data in enumerate(df.itertuples(index=False), start=start_row + 1):
        alt = (ri % 2 == 0)
        for ci, value in enumerate(row_data, start=1):
            cell = ws.cell(row=ri, column=ci, value=value)
            cell.font      = FONT_NORMAL
            cell.alignment = ALIGN_RIGHT

            col_name = df.columns[ci - 1]

            if alt:
                cell.fill = F_ALT

            # Formato numérico
            if col_name in COL_FORMATS:
                cell.number_format = COL_FORMATS[col_name]
            elif col_name in ("Periodo",):
                cell.alignment = ALIGN_LEFT
                cell.font = FONT_NORMAL
            elif isinstance(value, float):
                cell.number_format = FMT_CURR

            # Colorear resultado neto
            if col_name == "Resultado_neto" and isinstance(value, (int, float)):
                cell.fill = F_POS if value >= 0 else F_NEG
                cell.font = FONT_POS if value >= 0 else FONT_NEG

    # Congelar fila de cabecera
    if freeze:
        ws.freeze_panes = ws.cell(row=start_row + 1, column=1)

    # Ajustar anchos
    for ci, col in enumerate(df.columns, start=1):
        letter = get_column_letter(ci)
        header_len = len(COL_NAMES_ES.get(col, col))
        max_data   = max(
            (len(str(v)) for v in df[col] if v is not None),
            default=0,
        )
        ws.column_dimensions[letter].width = min(max(header_len, max_data) + 3, 22)

    # Altura de cabecera
    ws.row_dimensions[start_row].height = 32


def export_to_excel(df_monthly: pd.DataFrame, params: ProjectionParams) -> bytes:
    """Genera workbook Excel y retorna bytes para descarga."""
    wb = Workbook()

    # ── Hoja 1: Resumen Ejecutivo ─────────────────────────────────────────────
    ws_sum = wb.active
    ws_sum.title = "Resumen Ejecutivo"
    ws_sum.sheet_view.showGridLines = False
    ws_sum.column_dimensions["A"].width = 3
    ws_sum.column_dimensions["B"].width = 32
    ws_sum.column_dimensions["C"].width = 22

    ws_sum.merge_cells("B2:F2")
    ws_sum["B2"] = "PROYECCIÓN FINANCIERA — FACTORING"
    ws_sum["B2"].font      = FONT_TITLE
    ws_sum["B2"].fill      = F_DARK
    ws_sum["B2"].alignment = ALIGN_CENTER

    ws_sum["B3"] = f"Generado: {datetime.now().strftime('%d/%m/%Y %H:%M')}  |  Horizonte: 6 años (72 meses)"
    ws_sum["B3"].font = FONT_LABEL

    # Parámetros base
    ws_sum["B5"] = "PARÁMETROS BASE"
    ws_sum["B5"].font = FONT_SUB
    ws_sum["B5"].fill = F_NAVY

    param_rows = [
        ("Cartera Inicial",              f"${params.initial_portfolio:,.0f}M"),
        ("Crecimiento Anual de Cartera", f"${params.annual_portfolio_growth:,.0f}M"),
        ("Tasa de Colocación (mensual)", f"{params.placement_rate:.2f}%"),
        ("Costo de Fondo (mensual)",     f"{params.funding_cost_rate:.2f}%"),
        ("Costos Operacionales Base",    f"${params.current_op_costs:,.1f}M/mes"),
        ("Remuneraciones Base",          f"${params.current_remuneration:,.1f}M/mes"),
        ("Incremento Anual Op.",         f"${params.annual_op_increment:,.1f}M/año"),
        ("Otros Gastos",                 f"${params.other_expenses:,.1f}M/mes"),
    ]
    if params.overrides:
        param_rows.append(("Modificaciones Programadas", f"{len(params.overrides)} override(s)"))

    for i, (label, value) in enumerate(param_rows, start=6):
        ws_sum[f"B{i}"] = label
        ws_sum[f"C{i}"] = value
        ws_sum[f"B{i}"].font = FONT_LABEL
        ws_sum[f"C{i}"].font = Font(name="Calibri", color="FFFFFF", size=10, bold=True)
        ws_sum[f"B{i}"].alignment = ALIGN_LEFT
        ws_sum[f"C{i}"].alignment = ALIGN_LEFT
        if i % 2 == 0:
            ws_sum[f"B{i}"].fill = F_ALT
            ws_sum[f"C{i}"].fill = F_ALT

    # KPIs finales
    last = df_monthly.iloc[-1]
    total_net  = df_monthly["Resultado_neto"].sum()
    avg_margin = df_monthly["Margen_neto_pct"].mean()
    final_portfolio = last["Cartera"]

    kpi_start = 18
    ws_sum[f"B{kpi_start}"] = "KPIs PROYECTADOS (6 AÑOS)"
    ws_sum[f"B{kpi_start}"].font = FONT_SUB
    ws_sum[f"B{kpi_start}"].fill = F_NAVY

    kpi_rows = [
        ("Cartera Final (Mes 72)",        f"${final_portfolio:,.0f}M"),
        ("Resultado Neto Total",           f"${total_net:,.0f}M"),
        ("Margen Neto Promedio",           f"{avg_margin:.1f}%"),
        ("ROA Final (Mes 72)",             f"{last['ROA_pct']:.2f}%"),
        ("Eficiencia Final (Mes 72)",      f"{last['Eficiencia_pct']:.1f}%"),
    ]
    for i, (label, value) in enumerate(kpi_rows, start=kpi_start + 1):
        ws_sum[f"B{i}"] = label
        ws_sum[f"C{i}"] = value
        ws_sum[f"B{i}"].font = FONT_LABEL
        ws_sum[f"C{i}"].font = FONT_POS if "M" in str(value) or "%" in str(value) else FONT_NORMAL
        if i % 2 == 0:
            ws_sum[f"B{i}"].fill = F_ALT
            ws_sum[f"C{i}"].fill = F_ALT

    # Tabla anual en hoja resumen
    df_annual = aggregate_by_period(df_monthly, "Anual")
    sum_cols  = ["Periodo", "Cartera", "Ingresos_factoring", "Margen_financiero",
                 "Total_costos", "Resultado_neto", "Margen_neto_pct", "ROA_pct"]
    df_a_disp = df_annual[[c for c in sum_cols if c in df_annual.columns]]

    ws_sum[f"B{kpi_start + 9}"] = "RESUMEN ANUAL"
    ws_sum[f"B{kpi_start + 9}"].font = FONT_SUB
    ws_sum[f"B{kpi_start + 9}"].fill = F_NAVY

    _write_df(ws_sum, df_a_disp, start_row=kpi_start + 10, freeze=False)

    # ── Hojas de proyección ───────────────────────────────────────────────────
    periods = [
        ("Proyección Mensual",    "Mensual"),
        ("Trimestral",            "Trimestral"),
        ("Semestral",             "Semestral"),
        ("Anual",                 "Anual"),
    ]

    monthly_cols = ["Mes", "Año", "Cartera", "Ingresos_factoring", "Costo_fondo",
                    "Margen_financiero", "Costos_operacionales", "Remuneraciones",
                    "Otros_gastos", "Total_costos", "Resultado_neto",
                    "Margen_neto_pct", "Eficiencia_pct", "ROA_pct"]

    period_cols = ["Periodo", "Cartera", "Ingresos_factoring", "Costo_fondo",
                   "Margen_financiero", "Costos_operacionales", "Remuneraciones",
                   "Total_costos", "Resultado_neto", "Margen_neto_pct", "Eficiencia_pct", "ROA_pct"]

    for sheet_name, period in periods:
        ws = wb.create_sheet(sheet_name)
        df_p = aggregate_by_period(df_monthly, period)

        if period == "Mensual":
            cols = [c for c in monthly_cols if c in df_p.columns]
        else:
            cols = [c for c in period_cols if c in df_p.columns]

        _write_df(ws, df_p[cols])

    # ── Serializar a bytes ────────────────────────────────────────────────────
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.read()


# ── PowerPoint export ─────────────────────────────────────────────────────────

def export_to_pptx(df_monthly: pd.DataFrame, params: ProjectionParams,   # noqa: C901
                   scenario_name: str = "") -> bytes:
    """Genera presentación ejecutiva PowerPoint y retorna bytes para descarga."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.oxml.ns import qn
    from lxml import etree

    # ── Paleta SMB ────────────────────────────────────────────────────────────
    C_DARK    = RGBColor(0x0A, 0x14, 0x28)   # fondo base
    C_PANEL   = RGBColor(0x12, 0x24, 0x44)   # panel interior
    C_NAVY    = RGBColor(0x1E, 0x3A, 0x5F)   # panel header
    C_ACCENT  = RGBColor(0x0D, 0x2D, 0x52)   # fila alternada
    C_ORANGE  = RGBColor(0xF4, 0x79, 0x20)   # naranja SMB
    C_ORANGE2 = RGBColor(0xFF, 0xA9, 0x4D)   # naranja claro
    C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
    C_LGRAY   = RGBColor(0xCB, 0xD5, 0xE1)
    C_GRAY    = RGBColor(0x7A, 0x8A, 0x9A)
    C_GREEN   = RGBColor(0x4A, 0xDE, 0x80)
    C_RED     = RGBColor(0xF8, 0x71, 0x71)
    C_PURPLE  = RGBColor(0xA7, 0x8B, 0xFA)
    C_TEAL    = RGBColor(0x2D, 0xD4, 0xBF)

    import os as _os
    W = Inches(13.33)
    H = Inches(7.5)

    # Rutas de assets (logo y decorativo geométrico extraídos del PPTX de referencia)
    _here      = _os.path.dirname(_os.path.abspath(__file__))
    LOGO_PATH  = _os.path.join(_here, "assets", "logo_smb.png")
    DECO_PATH  = _os.path.join(_here, "assets", "deco.png")
    HAS_LOGO   = _os.path.exists(LOGO_PATH)
    HAS_DECO   = _os.path.exists(DECO_PATH)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    df_annual = aggregate_by_period(df_monthly, "Anual")
    last      = df_monthly.iloc[-1]
    m1        = df_monthly.iloc[0]
    total_net = df_monthly["Resultado_neto"].sum()
    ts        = datetime.now().strftime("%d/%m/%Y")
    footer_txt = f"SMB Servicios Financieros S.A.  ·  Proyección 6 años  ·  {ts}"

    # ── Helpers ───────────────────────────────────────────────────────────────

    def bg(slide, color: RGBColor = None):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color or C_DARK

    def rect(slide, l, t, w, h, color, radius=False):
        shp_type = 5 if radius else 1   # 5=ROUNDED_RECTANGLE
        s = slide.shapes.add_shape(shp_type, l, t, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        if radius:
            s.adjustments[0] = 0.05
        return s

    def txt(slide, text, l, t, w, h, size=11, bold=False,
            color=None, align=PP_ALIGN.LEFT, italic=False):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        p  = tf.paragraphs[0]
        p.alignment = align
        r  = p.add_run()
        r.text        = text
        r.font.size   = Pt(size)
        r.font.bold   = bold
        r.font.italic = italic
        r.font.color.rgb = color or C_WHITE
        return tb

    def slide_header(slide, title, subtitle="", slide_num=""):
        """Barra superior naranja + título + logo SMB."""
        rect(slide, 0, 0, W, Inches(0.07), C_ORANGE)
        rect(slide, 0, Inches(0.07), W, Inches(1.05), C_PANEL)
        # línea vertical naranja izquierda
        rect(slide, Inches(0.35), Inches(0.14), Inches(0.045), Inches(0.9), C_ORANGE)
        # Logo en header: fondo blanco redondeado + imagen con ratio correcto
        LOGO_W = Inches(2.15)
        LOGO_PAD_H = Inches(0.12)
        LOGO_PAD_V = Inches(0.1)
        logo_l = W - LOGO_W - Inches(0.18)
        logo_t = Inches(0.1)
        logo_box_h = Inches(0.9)
        _s = slide.shapes.add_shape(5, logo_l - LOGO_PAD_H, logo_t,
                                    LOGO_W + LOGO_PAD_H * 2, logo_box_h)
        _s.fill.solid(); _s.fill.fore_color.rgb = C_WHITE
        _s.line.fill.background()
        _s.adjustments[0] = 0.08
        add_logo(slide, logo_l, logo_t + LOGO_PAD_V, LOGO_W)
        # Número de slide (izq. del logo)
        if slide_num:
            txt(slide, slide_num, logo_l - Inches(0.85), Inches(0.2), Inches(0.65), Inches(0.45),
                size=18, bold=True, color=C_ORANGE, align=PP_ALIGN.RIGHT)
        # Título y subtítulo
        txt(slide, title, Inches(0.5), Inches(0.12), logo_l - Inches(1.0), Inches(0.62),
            size=24, bold=True, color=C_WHITE)
        if subtitle:
            txt(slide, subtitle, Inches(0.5), Inches(0.7), logo_l - Inches(1.0), Inches(0.38),
                size=10, color=C_GRAY)
        # footer
        rect(slide, 0, H - Inches(0.28), W, Inches(0.28), C_PANEL)
        rect(slide, 0, H - Inches(0.28), W, Inches(0.02), C_ORANGE)
        txt(slide, footer_txt, Inches(0.3), H - Inches(0.27), Inches(12), Inches(0.27),
            size=7, color=C_GRAY)

    def panel(slide, l, t, w, h, title, accent_color=None):
        """Panel con borde superior de color."""
        ac = accent_color or C_ORANGE
        rect(slide, l, t, w, h, C_PANEL, radius=False)
        rect(slide, l, t, w, Inches(0.04), ac)
        txt(slide, title, l + Inches(0.14), t + Inches(0.07),
            w - Inches(0.2), Inches(0.32), size=9, bold=True, color=ac)
        return t + Inches(0.42)   # y de inicio del contenido

    def item_row(slide, l, t, w, label, value, val_color=None, alt=False):
        """Fila etiqueta | valor dentro de un panel."""
        if alt:
            rect(slide, l, t, w, Inches(0.295), C_ACCENT)
        vc = val_color or C_LGRAY
        txt(slide, label, l + Inches(0.14), t + Inches(0.04),
            w * 0.58, Inches(0.27), size=9, color=C_GRAY)
        txt(slide, value, l + w * 0.58, t + Inches(0.04),
            w * 0.38, Inches(0.27), size=9, bold=True, color=vc,
            align=PP_ALIGN.RIGHT)
        return t + Inches(0.305)

    def kpi_card(slide, l, t, w, h, label, value, sub="", pos=True):
        """Tarjeta KPI con borde izquierdo naranja."""
        rect(slide, l, t, w, h, C_PANEL)
        rect(slide, l, t, Inches(0.045), h, C_ORANGE if pos else C_RED)
        txt(slide, label, l + Inches(0.1), t + Inches(0.08),
            w - Inches(0.15), Inches(0.32), size=8, color=C_GRAY, align=PP_ALIGN.CENTER)
        vclr = C_GREEN if pos else C_RED
        txt(slide, value, l + Inches(0.1), t + Inches(0.35),
            w - Inches(0.15), Inches(0.58), size=20, bold=True, color=vclr,
            align=PP_ALIGN.CENTER)
        if sub:
            txt(slide, sub, l + Inches(0.1), t + h - Inches(0.3),
                w - Inches(0.15), Inches(0.28), size=7, color=C_GRAY,
                align=PP_ALIGN.CENTER, italic=True)

    def add_logo(slide, l, t, w):
        """Inserta logo SMB preservando aspect ratio (solo ancho controlado)."""
        if HAS_LOGO:
            slide.shapes.add_picture(LOGO_PATH, l, t, width=w)

    def style_chart(chart, bg_color=None):
        """Fondo transparente y quitar borde al área de plot."""
        try:
            chart.plot_area.format.fill.background()
            chart.plot_area.format.line.fill.background()
            chart.chart_area.format.fill.background()
            chart.chart_area.format.line.fill.background()
        except Exception:
            pass

    def fix_invert_negative(chart):
        """Desactiva la inversión de color en barras negativas (val=0).

        El elemento <c:invertIfNegative> debe ir DESPUÉS de <c:spPr> en el
        schema OOXML — etree.SubElement lo pone al final lo que PowerPoint ignora.
        Usamos insert() en la posición correcta.
        """
        try:
            for plot in chart.plots:
                for ser in plot.series:
                    sp = ser._element
                    # Quitar si ya existe (posición incorrecta)
                    existing = sp.find(qn("c:invertIfNegative"))
                    if existing is not None:
                        sp.remove(existing)
                    # Insertar justo después de c:spPr
                    children = list(sp)
                    spPr_el = sp.find(qn("c:spPr"))
                    pos = (children.index(spPr_el) + 1) if spPr_el is not None else len(children)
                    inv = etree.Element(qn("c:invertIfNegative"))
                    inv.set("val", "0")
                    sp.insert(pos, inv)
        except Exception:
            pass

    def add_data_labels(chart, font_size=8, color="FFFFFF"):
        """Agrega etiquetas de datos a todos los series."""
        try:
            for plot in chart.plots:
                for ser in plot.series:
                    sp = ser._element
                    dLbls = etree.SubElement(sp, qn("c:dLbls"))
                    for tag, val in [("showVal", "1"), ("showLegendKey", "0"),
                                     ("showCatName", "0"), ("showSerName", "0"),
                                     ("showPercent", "0"), ("showBubbleSize", "0")]:
                        el = etree.SubElement(dLbls, qn(f"c:{tag}"))
                        el.set("val", val)
                    txPr = etree.SubElement(dLbls, qn("c:txPr"))
                    etree.SubElement(txPr, qn("a:bodyPr"))
                    etree.SubElement(txPr, qn("a:lstStyle"))
                    p_el  = etree.SubElement(txPr, qn("a:p"))
                    pPr   = etree.SubElement(p_el, qn("a:pPr"))
                    defRP = etree.SubElement(pPr, qn("a:defRPr"))
                    defRP.set("sz", str(int(font_size * 100)))
                    defRP.set("b", "1")
                    sf   = etree.SubElement(defRP, qn("a:solidFill"))
                    clr  = etree.SubElement(sf, qn("a:srgbClr"))
                    clr.set("val", color)
        except Exception:
            pass

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 1 — PORTADA
    # ══════════════════════════════════════════════════════════════════════════
    s1 = prs.slides.add_slide(blank)
    bg(s1, C_DARK)

    # Decorativo geométrico en la derecha (mismo que PPTX de referencia)
    if HAS_DECO:
        s1.shapes.add_picture(DECO_PATH, W - Inches(3.8), 0, width=Inches(3.8), height=H)

    # Barras naranja izquierda
    rect(s1, 0, 0, Inches(0.55), H, C_ORANGE)
    rect(s1, Inches(0.55), 0, Inches(0.12), H, C_ORANGE2)

    # Logo SMB en portada: caja blanca + imagen con ratio correcto
    _LPAD = Inches(0.15)
    _LBOX_W = Inches(4.1)
    _LBOX_H = Inches(1.5)
    _s_logo = s1.shapes.add_shape(5, Inches(0.85), Inches(0.45), _LBOX_W, _LBOX_H)
    _s_logo.fill.solid(); _s_logo.fill.fore_color.rgb = C_WHITE
    _s_logo.line.fill.background(); _s_logo.adjustments[0] = 0.06
    add_logo(s1, Inches(0.85) + _LPAD, Inches(0.45) + _LPAD, _LBOX_W - _LPAD * 2)

    # Línea divisoria naranja bajo el logo (logo termina en 0.45+1.5=1.95")
    rect(s1, Inches(0.85), Inches(2.08), Inches(8.8), Inches(0.035), C_ORANGE)

    title = scenario_name or "Proyección Financiera"
    txt(s1, title,
        Inches(0.85), Inches(2.2), Inches(9.0), Inches(1.7),
        size=40, bold=True, color=C_WHITE)

    txt(s1, "Factoring  +  Leasing  |  Horizonte 6 años (72 meses)",
        Inches(0.85), Inches(3.98), Inches(9.0), Inches(0.5),
        size=14, color=C_LGRAY)
    txt(s1, "Valores expresados en millones de pesos chilenos (M$)",
        Inches(0.85), Inches(4.52), Inches(9.0), Inches(0.4),
        size=11, color=C_GRAY, italic=True)

    # Caja info fecha
    rect(s1, Inches(0.85), Inches(5.4), Inches(3.8), Inches(1.55), C_PANEL, radius=True)
    txt(s1, "FECHA DE GENERACIÓN", Inches(1.0), Inches(5.55), Inches(3.4), Inches(0.35),
        size=7, color=C_GRAY, align=PP_ALIGN.CENTER)
    txt(s1, ts, Inches(1.0), Inches(5.88), Inches(3.4), Inches(0.55),
        size=20, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)
    txt(s1, "Calibrada con EERR Mayo 2026",
        Inches(1.0), Inches(6.42), Inches(3.4), Inches(0.4),
        size=8, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)

    # Footer portada
    rect(s1, 0, H - Inches(0.35), W, Inches(0.35), RGBColor(0x06, 0x0E, 0x1A))
    txt(s1, footer_txt, Inches(0.7), H - Inches(0.33), Inches(12), Inches(0.33),
        size=7, color=C_GRAY)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 2 — SUPUESTOS: TASAS Y CARTERAS
    # ══════════════════════════════════════════════════════════════════════════
    s2 = prs.slides.add_slide(blank)
    bg(s2)
    slide_header(s2, "SUPUESTOS DEL MODELO",
                 "Parámetros financieros y carteras iniciales — calibrado EERR Mayo 2026", "2")

    # Calcular spread y tasa lineal leasing
    spread = params.placement_rate - params.funding_cost_rate
    r_les  = params.leasing_annual_rate / 12 / 100
    n_les  = params.leasing_avg_term_months
    if r_les > 0:
        cuota = r_les / (1 - (1 + r_les) ** (-n_les))
        tasa_lineal = (cuota * n_les - 1) / n_les * 100
    else:
        tasa_lineal = 0.0
    leas_mm = params.leasing_portfolio_uf * params.uf_value / 1_000_000

    CW = Inches(4.1)
    GAP = Inches(0.2)
    ML  = Inches(0.3)
    TY  = Inches(1.22)
    PH  = Inches(5.98)

    # ── Panel 1: Tasas y Financiamiento ──────────────────────────────────────
    p1x = ML
    cy = panel(s2, p1x, TY, CW, PH, "TASAS Y FINANCIAMIENTO", C_ORANGE)
    rows1 = [
        ("Tasa Colocación Factoring",  f"{params.placement_rate:.2f}%/mes", C_ORANGE),
        ("Equivalente anual (ef.)",    f"{((1+params.placement_rate/100)**12-1)*100:.1f}%/año", None),
        ("Tasa Interés Leasing (TIR)", f"{params.leasing_annual_rate:.1f}%/año", C_ORANGE),
        ("Ingr. lineal leasing",       f"{tasa_lineal:.3f}%/mes", None),
        ("Costo de Fondo (variable)",  f"{params.funding_cost_rate:.2f}%/mes", C_RED),
        ("Equivalente anual",          f"{params.funding_cost_rate*12:.2f}%/año", None),
        ("Costo Fijo de Fondo",        f"M${params.base_funding_cost_mm:.2f}/mes", C_RED),
        ("Spread Neto Factoring",      f"{spread:.2f}%/mes", C_GREEN if spread > 0 else C_RED),
        ("IPC Proyectado (reaj. UF)",  f"{params.monthly_ipc:.2f}%/mes", None),
        ("IPC anual equivalente",      f"{((1+params.monthly_ipc/100)**12-1)*100:.1f}%/año", None),
        ("Plazo Prom. Leasing",        f"{params.leasing_avg_term_months} meses", None),
        ("Impuesto Renta",             f"{params.tax_rate:.0f}%", None),
    ]
    for i, (lbl, val, vc) in enumerate(rows1):
        cy = item_row(s2, p1x, cy, CW, lbl, val, val_color=vc, alt=(i % 2 == 0))

    # ── Panel 2: Carteras Iniciales ──────────────────────────────────────────
    p2x = ML + CW + GAP
    cy = panel(s2, p2x, TY, CW, PH, "CARTERAS INICIALES (EERR MAY-26)", C_TEAL)
    rows2 = [
        ("Cartera Factoring",          f"M${params.initial_portfolio:,.0f}", C_ORANGE),
        ("Cartera Leasing (UF)",       f"{params.leasing_portfolio_uf:,.0f} UF", C_PURPLE),
        ("Cartera Leasing (M$)",       f"M${leas_mm:,.0f}", C_PURPLE),
        ("Cartera Total Inicial",      f"M${params.initial_portfolio+leas_mm:,.0f}", C_TEAL),
        ("Valor UF Base",              f"${params.uf_value:,.0f}", None),
        ("Crecim. Factoring/año",      f"M${params.annual_portfolio_growth:,.0f}/año", C_ORANGE),
        ("Crecim. Factoring/mes",      f"M${params.annual_portfolio_growth/12:,.1f}/mes", None),
        ("Crecim. Leasing/año",        f"M${params.leasing_annual_growth_mm:,.0f}/año", C_PURPLE),
        ("Crecim. Leasing/mes",        f"M${params.leasing_annual_growth_mm/12:,.1f}/mes", None),
        ("Ingr. Factoring Mes 1",      f"M${m1['Ingresos_factoring']:,.1f}", C_ORANGE),
        ("Ingr. Leasing Mes 1",        f"M${m1['Ingresos_leasing']:,.1f}", C_PURPLE),
        ("Ingr. Total Mes 1",          f"M${m1['Ingresos_total']:,.1f}", C_TEAL),
    ]
    for i, (lbl, val, vc) in enumerate(rows2):
        cy = item_row(s2, p2x, cy, CW, lbl, val, val_color=vc, alt=(i % 2 == 0))

    # ── Panel 3: Costos y Riesgo ─────────────────────────────────────────────
    p3x = ML + 2 * (CW + GAP)
    cy = panel(s2, p3x, TY, CW, PH, "COSTOS OPERACIONALES Y RIESGO", C_ORANGE2)
    total_cos_mes = params.current_remuneration + params.current_op_costs + params.other_expenses
    rows3 = [
        ("Remuneraciones Base",        f"M${params.current_remuneration:,.1f}/mes", C_LGRAY),
        ("Remuneraciones / año",       f"M${params.current_remuneration*12:,.0f}/año", None),
        ("Costos Operacionales Base",  f"M${params.current_op_costs:,.1f}/mes", C_LGRAY),
        ("Otros Gastos No Op.",        f"M${params.other_expenses:,.1f}/mes", C_LGRAY),
        ("TOTAL COSTOS FIJOS/mes",     f"M${total_cos_mes:,.1f}/mes", C_RED),
        ("TOTAL COSTOS FIJOS/año",     f"M${total_cos_mes*12:,.0f}/año", C_RED),
        ("Incremento Anual Costos Op.",f"M${params.annual_op_increment:,.1f}/año", None),
        ("Tasa Provisiones",           f"{params.provision_rate:.1f}%/año cartera", C_ORANGE),
        ("Provisión Mes 1",            f"M${m1['Provision_expense']:,.1f}/mes", None),
        ("NPL Estimado",               f"{params.npl_rate:.1f}% cartera morosa", None),
        ("Patrimonio Inicial",         f"M${params.initial_equity:,.0f}", None),
        ("Modificaciones Programa.",   f"{len(params.overrides)} tramos", C_ORANGE if params.overrides else None),
    ]
    for i, (lbl, val, vc) in enumerate(rows3):
        cy = item_row(s2, p3x, cy, CW, lbl, val, val_color=vc, alt=(i % 2 == 0))

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 3 — MODIFICACIONES PROGRAMADAS (solo si hay overrides)
    # ══════════════════════════════════════════════════════════════════════════
    growth_ovs   = [o for o in params.overrides if o.param == "portfolio_growth"]
    personal_ovs = [o for o in params.overrides if o.param in ("remuneration", "operational_costs", "otros_gastos")]
    other_ovs    = [o for o in params.overrides if o not in growth_ovs and o not in personal_ovs]

    if params.overrides:
        s3 = prs.slides.add_slide(blank)
        bg(s3)
        slide_header(s3, "MODIFICACIONES PROGRAMADAS",
                     f"{len(params.overrides)} cambio(s) de parámetros a lo largo del horizonte", "3")

        LABEL_MAP = {
            "portfolio_growth":  "Crecimiento Factoring",
            "leasing_growth":    "Crecimiento Leasing",
            "funding_cost":      "Costo de Fondo (%)",
            "base_funding_cost": "Costo Fijo Fondo (M$)",
            "placement_rate":    "Tasa Colocación (%)",
            "remuneration":      "Remuneraciones",
            "operational_costs": "Costos Operacionales",
            "otros_gastos":      "Otros Gastos",
            "leasing_rate":      "Tasa Leasing",
            "provision_rate":    "Tasa Provisiones",
            "tax_rate":          "Impuesto Renta",
            "ipc_mensual":       "IPC Mensual",
        }
        UNIT_MAP = {
            "portfolio_growth": "M$/mes", "leasing_growth": "M$/año",
            "funding_cost": "%/mes", "base_funding_cost": "M$/mes",
            "placement_rate": "%/mes", "remuneration": "M$/mes",
            "operational_costs": "M$/mes", "otros_gastos": "M$/mes",
            "leasing_rate": "%/año", "provision_rate": "%/año",
            "tax_rate": "%", "ipc_mensual": "%/mes",
        }

        groups = [
            ("CRECIMIENTO DE CARTERA",  growth_ovs,   C_ORANGE),
            ("PERSONAL Y COSTOS",       personal_ovs, C_ORANGE2),
            ("OTROS PARÁMETROS",        other_ovs,    C_TEAL),
        ]
        GW = Inches(4.1)
        gx = ML
        for g_title, g_ovs, g_clr in groups:
            gy = panel(s3, gx, TY, GW, PH, g_title, g_clr)
            if not g_ovs:
                txt(s3, "Sin modificaciones en este grupo",
                    gx + Inches(0.2), gy + Inches(0.1), GW - Inches(0.3), Inches(0.35),
                    size=9, color=C_GRAY, italic=True)
            else:
                for i, ov in enumerate(sorted(g_ovs, key=lambda x: x.from_month)):
                    yr  = (ov.from_month - 1) // 12 + 1
                    lbl = LABEL_MAP.get(ov.param, ov.param)
                    unit = UNIT_MAP.get(ov.param, "")
                    typ  = "+" if ov.is_delta else "→"
                    val_str = f"{typ} {ov.value:,.1f} {unit}"
                    mes_str = f"Mes {ov.from_month} (Año {yr})"
                    if ov.note:
                        mes_str += f"  ·  {ov.note}"
                    gy = item_row(s3, gx, gy, GW, f"{lbl}  {mes_str}", val_str,
                                  val_color=g_clr, alt=(i % 2 == 0))
            gx += GW + GAP
        slide_num_start = 4
    else:
        slide_num_start = 3

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE KPIs
    # ══════════════════════════════════════════════════════════════════════════
    sk = prs.slides.add_slide(blank)
    bg(sk)
    slide_header(sk, "INDICADORES CLAVE DE LA PROYECCIÓN",
                 "Resultado acumulado a 6 años  ·  Evolución Mes 1 → Mes 72",
                 str(slide_num_start))

    KW  = Inches(2.42)
    KH  = Inches(1.42)
    KG  = Inches(0.155)
    KT1 = Inches(1.22)
    KT2 = KT1 + KH + KG
    KT3 = KT2 + KH + KG

    kpis = [
        # (label, value, sublabel, pos, row, col)
        ("Cartera Total Final",        f"M${last['Cartera_total']:,.0f}",
         "Factoring + Leasing · Mes 72",    True, 0, 0),
        ("Cartera Factoring Final",    f"M${last['Cartera']:,.0f}",
         "Mes 72",                           True, 0, 1),
        ("Cartera Leasing Final",      f"M${last['Cartera_leasing']:,.0f}",
         "Mes 72",                           True, 0, 2),
        ("Ingresos Totales Año 6",     f"M${df_annual.iloc[-1]['Ingresos_total']:,.0f}",
         "Factoring + Leasing",             True, 0, 3),
        ("Resultado Neto Acumulado",   f"M${total_net:,.0f}",
         "6 años consolidado",              total_net >= 0, 0, 4),
        ("Resultado Neto Mes 1",       f"M${m1['Resultado_neto']:,.1f}",
         "Base EERR May-26",               m1["Resultado_neto"] >= 0, 1, 0),
        ("Resultado Neto Mes 72",      f"M${last['Resultado_neto']:,.1f}",
         "Estado estable",                  last["Resultado_neto"] >= 0, 1, 1),
        ("Margen Neto Promedio",       f"{df_monthly['Margen_neto_pct'].mean():.1f}%",
         "Prom. 72 meses",                  True, 1, 2),
        ("ROA Anualizado (Mes 72)",    f"{last['ROA_pct']:.2f}%",
         "Retorno sobre activos",           last['ROA_pct'] >= 0, 1, 3),
        ("Eficiencia (Mes 72)",        f"{last['Eficiencia_pct']:.1f}%",
         "Costos / Margen financiero",      last['Eficiencia_pct'] < 100, 1, 4),
        ("Patrimonio Final",           f"M${last['Patrimonio']:,.0f}",
         "Mes 72 acumulado",               True, 2, 0),
        ("Costo Fondo Total/año 6",    f"M${df_annual.iloc[-1]['Costo_fondo']:,.0f}",
         "Año 6 anualizado",               True, 2, 1),
        ("Margen Bruto Año 6",         f"M${df_annual.iloc[-1]['Margen_financiero']:,.0f}",
         "Ingresos − Costo de Fondo",      True, 2, 2),
        ("Total Costos Año 6",         f"M${abs(df_annual.iloc[-1]['Total_costos']):,.0f}",
         "Op. + Rem. + Otros",             True, 2, 3),
        ("Impuesto Pagado (6 años)",   f"M${df_monthly['Impuesto'].sum():,.0f}",
         "27% sobre resultado positivo",   True, 2, 4),
    ]

    KT_map = {0: KT1, 1: KT2, 2: KT3}
    for lbl, val, sub, pos, row, col in kpis:
        kpi_card(sk, ML + col * (KW + KG), KT_map[row], KW, KH, lbl, val, sub, pos)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE — TABLA ANUAL P&L
    # ══════════════════════════════════════════════════════════════════════════
    st4 = prs.slides.add_slide(blank)
    bg(st4)
    slide_header(st4, "RESUMEN ANUAL  —  ESTADO DE RESULTADOS",
                 "Valores en M$ · Factoring + Leasing consolidado",
                 str(slide_num_start + 1))

    tbl_cols = ["Periodo", "Ingresos_factoring", "Ingresos_leasing",
                "Ingresos_total", "Costo_fondo", "Margen_financiero",
                "Total_costos", "EBIT", "Resultado_neto", "Margen_neto_pct"]
    tbl_hdrs = ["Período", "Ing.Fact.", "Ing.Leas.", "Ing.Total",
                "C.Fondo", "Margen", "Tot.Costos", "EBIT", "Res.Neto", "M.Neto%"]

    n_rows = len(df_annual) + 1
    n_cols = len(tbl_cols)
    tbl = st4.shapes.add_table(
        n_rows, n_cols,
        Inches(0.25), Inches(1.22),
        W - Inches(0.5), H - Inches(1.5)
    ).table

    col_ws = [Inches(1.25)] + [Inches(1.23)] * (n_cols - 1)
    for ci, cw in enumerate(col_ws):
        tbl.columns[ci].width = cw
    tbl.rows[0].height = int(Inches(0.42))

    for ci, hdr in enumerate(tbl_hdrs):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x06, 0x10, 0x22)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = hdr
        r = p.runs[0]
        r.font.bold = True
        r.font.size = Pt(8.5)
        r.font.color.rgb = C_ORANGE

    for ri, (_, row) in enumerate(df_annual.iterrows(), start=1):
        bg_c = C_ACCENT if ri % 2 == 0 else C_PANEL
        tbl.rows[ri].height = int(Inches(0.48))
        for ci, col in enumerate(tbl_cols):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            val  = row[col]
            if col == "Periodo":
                txt_v = str(val);  clr = C_WHITE;  al = PP_ALIGN.LEFT
            elif col == "Margen_neto_pct":
                txt_v = f"{val:.1f}%";  clr = C_TEAL;  al = PP_ALIGN.CENTER
            else:
                txt_v = f"${val:,.0f}"
                if col == "Resultado_neto":
                    clr = C_GREEN if val >= 0 else C_RED
                elif col == "EBIT":
                    clr = C_LGRAY if val >= 0 else C_RED
                elif col == "Costo_fondo":
                    clr = C_RED
                elif col in ("Ingresos_factoring", "Ingresos_leasing", "Ingresos_total"):
                    clr = C_ORANGE
                elif col == "Margen_financiero":
                    clr = C_TEAL
                else:
                    clr = C_LGRAY
                al = PP_ALIGN.RIGHT
            cell.fill.fore_color.rgb = bg_c
            p = cell.text_frame.paragraphs[0]
            p.alignment = al
            p.text = txt_v
            r = p.runs[0]
            r.font.size = Pt(9)
            r.font.color.rgb = clr
            r.font.bold = col in ("Resultado_neto", "Ingresos_total", "Margen_financiero")

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE — GRÁFICO: RESULTADO NETO POR AÑO
    # ══════════════════════════════════════════════════════════════════════════
    s5 = prs.slides.add_slide(blank)
    bg(s5)
    slide_header(s5, "RESULTADO NETO ANUAL",
                 "Evolución del resultado neto proyectado año a año (M$)",
                 str(slide_num_start + 2))

    cd5 = ChartData()
    cd5.categories = [str(r["Periodo"]) for _, r in df_annual.iterrows()]
    vals5 = [round(r["Resultado_neto"], 1) for _, r in df_annual.iterrows()]
    cd5.add_series("Resultado Neto (M$)", vals5)

    chart5 = s5.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.18), Inches(12.3), Inches(5.95), cd5
    ).chart
    chart5.has_legend = False
    chart5.has_title  = False
    style_chart(chart5)
    chart5.plots[0].series[0].format.fill.solid()
    chart5.plots[0].series[0].format.fill.fore_color.rgb = C_ORANGE
    fix_invert_negative(chart5)
    add_data_labels(chart5, font_size=9)

    # Breakeven line (zero)
    try:
        va5 = chart5.value_axis
        va5.has_major_gridlines = True
        va5.major_gridlines.format.line.color.rgb = RGBColor(0x1E, 0x2E, 0x42)
        va5.tick_labels.font.size = Pt(9)
        va5.tick_labels.font.color.rgb = C_GRAY
        chart5.category_axis.tick_labels.font.size = Pt(9)
        chart5.category_axis.tick_labels.font.color.rgb = C_LGRAY
    except Exception:
        pass

    # Mini tabla resumen al costado
    mn_l = Inches(9.8)
    my   = Inches(1.25)
    rect(s5, mn_l, my, Inches(3.3), Inches(5.9), C_PANEL)
    txt(s5, "RESUMEN", mn_l + Inches(0.1), my + Inches(0.08),
        Inches(3.0), Inches(0.3), size=8, bold=True, color=C_ORANGE)
    for i, (_, row) in enumerate(df_annual.iterrows()):
        rty  = my + Inches(0.42) + i * Inches(0.73)
        bg_r = C_ACCENT if i % 2 == 0 else C_PANEL
        rect(s5, mn_l, rty, Inches(3.3), Inches(0.7), bg_r)
        rect(s5, mn_l, rty, Inches(0.04), Inches(0.7),
             C_GREEN if row["Resultado_neto"] >= 0 else C_RED)
        txt(s5, str(row["Periodo"]), mn_l + Inches(0.12), rty + Inches(0.04),
            Inches(1.6), Inches(0.3), size=8, color=C_LGRAY)
        clr_r = C_GREEN if row["Resultado_neto"] >= 0 else C_RED
        txt(s5, f"M${row['Resultado_neto']:,.0f}", mn_l + Inches(0.12), rty + Inches(0.34),
            Inches(3.0), Inches(0.3), size=11, bold=True, color=clr_r)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE — GRÁFICO: EVOLUCIÓN DE CARTERA
    # ══════════════════════════════════════════════════════════════════════════
    s6 = prs.slides.add_slide(blank)
    bg(s6)
    slide_header(s6, "EVOLUCIÓN DE CARTERA",
                 "Cartera de Factoring + Leasing proyectada año a año (M$)",
                 str(slide_num_start + 3))

    cd6 = ChartData()
    cd6.categories = [str(r["Periodo"]) for _, r in df_annual.iterrows()]
    cd6.add_series("Factoring", [round(r["Cartera"], 0) for _, r in df_annual.iterrows()])
    cd6.add_series("Leasing",   [round(r["Cartera_leasing"], 0) for _, r in df_annual.iterrows()])

    chart6 = s6.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        Inches(0.4), Inches(1.18), Inches(8.8), Inches(5.95), cd6
    ).chart
    chart6.has_legend = True
    chart6.has_title  = False
    style_chart(chart6)
    chart6.plots[0].series[0].format.fill.solid()
    chart6.plots[0].series[0].format.fill.fore_color.rgb = C_ORANGE
    chart6.plots[0].series[1].format.fill.solid()
    chart6.plots[0].series[1].format.fill.fore_color.rgb = C_PURPLE
    fix_invert_negative(chart6)
    try:
        chart6.legend.font.size = Pt(9)
        chart6.legend.font.color.rgb = C_LGRAY
        chart6.value_axis.tick_labels.font.color.rgb = C_GRAY
        chart6.value_axis.tick_labels.font.size = Pt(8)
        chart6.category_axis.tick_labels.font.color.rgb = C_LGRAY
        chart6.category_axis.tick_labels.font.size = Pt(9)
    except Exception:
        pass

    # Panel lateral: crecimiento anual
    px6 = Inches(9.4)
    py6 = Inches(1.22)
    rect(s6, px6, py6, Inches(3.7), Inches(5.9), C_PANEL)
    txt(s6, "CRECIMIENTO POR AÑO", px6 + Inches(0.1), py6 + Inches(0.08),
        Inches(3.4), Inches(0.3), size=8, bold=True, color=C_TEAL)
    prev_total = params.initial_portfolio + leas_mm
    for i, (_, row) in enumerate(df_annual.iterrows()):
        curr_total = row["Cartera_total"]
        incr = curr_total - prev_total
        rty  = py6 + Inches(0.42) + i * Inches(0.73)
        rect(s6, px6, rty, Inches(3.7), Inches(0.7),
             C_ACCENT if i % 2 == 0 else C_PANEL)
        txt(s6, str(row["Periodo"]), px6 + Inches(0.1), rty + Inches(0.04),
            Inches(2.0), Inches(0.3), size=8, color=C_LGRAY)
        txt(s6, f"M${curr_total:,.0f}", px6 + Inches(0.1), rty + Inches(0.34),
            Inches(2.2), Inches(0.3), size=10, bold=True, color=C_WHITE)
        txt(s6, f"+M${incr:,.0f}", px6 + Inches(2.3), rty + Inches(0.34),
            Inches(1.2), Inches(0.3), size=9, bold=True, color=C_ORANGE)
        prev_total = curr_total

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE — GRÁFICO: INGRESOS, COSTO FONDO Y MARGEN
    # ══════════════════════════════════════════════════════════════════════════
    s7 = prs.slides.add_slide(blank)
    bg(s7)
    slide_header(s7, "INGRESOS  ·  COSTO DE FONDO  ·  MARGEN",
                 "Factoring + Leasing · Evolución anual (M$)",
                 str(slide_num_start + 4))

    cd7 = ChartData()
    cd7.categories = [str(r["Periodo"]) for _, r in df_annual.iterrows()]
    cd7.add_series("Ing. Factoring",   [round(r["Ingresos_factoring"], 0) for _, r in df_annual.iterrows()])
    cd7.add_series("Ing. Leasing",     [round(r["Ingresos_leasing"], 0) for _, r in df_annual.iterrows()])
    cd7.add_series("Costo de Fondo",   [round(abs(r["Costo_fondo"]), 0) for _, r in df_annual.iterrows()])
    cd7.add_series("Margen Financiero",[round(r["Margen_financiero"], 0) for _, r in df_annual.iterrows()])
    cd7.add_series("Total Costos",     [round(abs(r["Total_costos"]), 0) for _, r in df_annual.iterrows()])

    chart7 = s7.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(0.4), Inches(1.18), Inches(9.0), Inches(5.95), cd7
    ).chart
    chart7.has_legend = True
    chart7.has_title  = False
    style_chart(chart7)
    clrs7 = [C_ORANGE, C_PURPLE, C_RED, C_TEAL, RGBColor(0xFB, 0xBF, 0x24)]
    for si, clr in enumerate(clrs7):
        try:
            chart7.plots[0].series[si].format.line.color.rgb = clr
            chart7.plots[0].series[si].format.line.width = Pt(2.5)
        except Exception:
            pass
    try:
        chart7.legend.font.size = Pt(9)
        chart7.legend.font.color.rgb = C_LGRAY
        chart7.value_axis.tick_labels.font.color.rgb = C_GRAY
        chart7.value_axis.tick_labels.font.size = Pt(8)
        chart7.value_axis.has_major_gridlines = True
        chart7.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0x1A, 0x2A, 0x3E)
        chart7.category_axis.tick_labels.font.color.rgb = C_LGRAY
        chart7.category_axis.tick_labels.font.size = Pt(9)
    except Exception:
        pass

    # Leyenda de colores manual
    lx7 = Inches(9.6)
    ly7 = Inches(1.3)
    rect(s7, lx7, ly7, Inches(3.5), Inches(5.8), C_PANEL)
    txt(s7, "VALORES AÑO A AÑO (M$)", lx7 + Inches(0.1), ly7 + Inches(0.1),
        Inches(3.2), Inches(0.3), size=8, bold=True, color=C_ORANGE)

    series_info = [
        ("Ing. Factoring",    "Ingresos_factoring",  C_ORANGE),
        ("Ing. Leasing",      "Ingresos_leasing",    C_PURPLE),
        ("Costo de Fondo",    "Costo_fondo",         C_RED),
        ("Margen Financiero", "Margen_financiero",   C_TEAL),
        ("Total Costos",      "Total_costos",        RGBColor(0xFB, 0xBF, 0x24)),
    ]
    ry = ly7 + Inches(0.48)
    for si_name, si_col, si_clr in series_info:
        txt(s7, si_name, lx7 + Inches(0.25), ry, Inches(1.8), Inches(0.28), size=8, color=C_LGRAY)
        rect(s7, lx7 + Inches(0.1), ry + Inches(0.06), Inches(0.1), Inches(0.14), si_clr)
        for ai, (_, arow) in enumerate(df_annual.iterrows()):
            val_s = arow[si_col]
            txt(s7, f"{val_s:,.0f}",
                lx7 + Inches(0.25) + Inches(0.52) * ai, ry + Inches(0.24),
                Inches(0.52), Inches(0.26),
                size=7, bold=True, color=si_clr, align=PP_ALIGN.RIGHT)
        ry += Inches(0.85)
        if ry < ly7 + Inches(5.7):
            rect(s7, lx7 + Inches(0.05), ry - Inches(0.06),
                 Inches(3.4), Inches(0.01), RGBColor(0x1A, 0x2A, 0x3E))

    # ── Serializar ────────────────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
